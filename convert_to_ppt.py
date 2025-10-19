import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def parse_markdown_slides(file_path):
    """Parse the markdown file and extract slide content"""
    with open(file_path, 'r', encoding='utf-8') as file:
        content = file.read()
    
    # Split by slide markers
    slides = re.split(r'---\n\n## Slide \d+:', content)
    
    # Remove the first empty split and clean up
    slides = [slide.strip() for slide in slides[1:] if slide.strip()]
    
    parsed_slides = []
    
    for slide in slides:
        # Extract title and content
        lines = slide.split('\n')
        
        # Find the title (first line with **)
        title = ""
        content_lines = []
        
        for line in lines:
            if line.startswith('**') and line.endswith('**'):
                title = line.replace('**', '').strip()
            elif line.strip():
                content_lines.append(line.strip())
        
        parsed_slides.append({
            'title': title,
            'content': content_lines
        })
    
    return parsed_slides

def is_english_text(text):
    """Check if text contains primarily English characters"""
    # Remove emojis and special characters first
    clean_text = re.sub(r'[^\w\s]', '', text)
    # Count English letters vs Chinese characters
    english_chars = len(re.findall(r'[a-zA-Z]', clean_text))
    chinese_chars = len(re.findall(r'[\u4e00-\u9fff]', clean_text))
    return english_chars > chinese_chars and english_chars > 0

def create_presentation(slides_data, output_file):
    """Create PowerPoint presentation from parsed slides"""
    prs = Presentation()
    
    # Use a layout that has title and content
    title_slide_layout = prs.slide_layouts[0]  # Title slide
    content_layout = prs.slide_layouts[1]  # Title and content
    
    for i, slide_data in enumerate(slides_data):
        if i == 0:
            # First slide - use title layout
            slide = prs.slides.add_slide(title_slide_layout)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            
            # Handle first slide specially
            title.text = "CapitalEpoch"
            subtitle.text = "GameFi DeFi Platform on Aptos\n\nBuilding Digital Commercial Streets Together"
            
        else:
            # Other slides - use content layout
            slide = prs.slides.add_slide(content_layout)
            title = slide.shapes.title
            content = slide.placeholders[1]
            
            title.text = slide_data['title']
            
            # Process content - only keep English content
            content_text = []
            for line in slide_data['content']:
                # Skip empty lines and slide markers
                if not line or 'Slide' in line and ':' in line:
                    continue
                
                # Skip Chinese-only lines
                if not is_english_text(line):
                    continue
                    
                # Handle different content types
                if line.startswith('❌') or line.startswith('✅') or line.startswith('🎮') or line.startswith('💰') or line.startswith('🤝') or line.startswith('🏆') or line.startswith('⚡') or line.startswith('📝') or line.startswith('⚛️') or line.startswith('🔗') or line.startswith('🤖') or line.startswith('🏪') or line.startswith('📊') or line.startswith('💱') or line.startswith('🎯') or line.startswith('📈') or line.startswith('🌐') or line.startswith('🔄') or line.startswith('🚀') or line.startswith('📈') or line.startswith('🌟') or line.startswith('🌍') or line.startswith('👨‍💼') or line.startswith('🏗️') or line.startswith('📄') or line.startswith('🌐') or line.startswith('📧') or line.startswith('🐦') or line.startswith('💬'):
                    # Bullet points with emojis
                    clean_line = line.replace('❌', '').replace('✅', '').replace('🎮', '').replace('💰', '').replace('🤝', '').replace('🏆', '').replace('⚡', '').replace('📝', '').replace('⚛️', '').replace('🔗', '').replace('🤖', '').replace('🏪', '').replace('📊', '').replace('💱', '').replace('🎯', '').replace('📈', '').replace('🌐', '').replace('🔄', '').replace('🚀', '').replace('🌟', '').replace('🌍', '').replace('👨‍💼', '').replace('🏗️', '').replace('📄', '').replace('📧', '').replace('🐦', '').replace('💬', '').strip()
                    # Only add if the cleaned line is English
                    if is_english_text(clean_line):
                        content_text.append(f"• {clean_line}")
                elif line.startswith('**') and line.endswith('**'):
                    # Bold text
                    clean_line = line.replace('**', '')
                    if is_english_text(clean_line):
                        content_text.append(clean_line)
                elif re.match(r'^\d+\.', line):
                    # Numbered list
                    if is_english_text(line):
                        content_text.append(f"  {line}")
                elif line.startswith('*') and line.endswith('*'):
                    # Italic text (usually translations) - skip these
                    continue
                else:
                    # Regular text
                    if is_english_text(line):
                        content_text.append(line)
            
            content.text = '\n'.join(content_text)
            
            # Format the content
            for paragraph in content.text_frame.paragraphs:
                paragraph.font.size = Pt(18)
                paragraph.font.name = 'Arial'
    
    # Save the presentation
    prs.save(output_file)
    print(f"Presentation saved as {output_file}")

def main():
    # Parse the English markdown file
    slides_data = parse_markdown_slides('PRESENTATION_ENGLISH.md')
    
    # Create PowerPoint presentation from English version
    create_presentation(slides_data, 'CapitalEpoch_Presentation_English.pptx')
    
    print(f"Successfully converted {len(slides_data)} slides to PowerPoint format from English-only markdown!")

if __name__ == "__main__":
    main()
